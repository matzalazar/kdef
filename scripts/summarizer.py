"""
summarizer.py — Generación de resúmenes markdown usando LLM

Responsabilidad única: tomar un archivo (PDF, DOCX, etc.), extraer su texto,
y generar un resumen estructurado en markdown usando un LLM.
Además, le pide al modelo un bloque machine-readable con fechas importantes
para que el pipeline pueda agregarlas a un calendario.

Decisiones de diseño:
  - La selección de proveedor LLM se hace por prefijo del nombre del modelo:
    "github/" usa la API compatible con OpenAI de GitHub Models
    "gemini/" usa la API nativa de Google Gemini
    Esto permite cambiar el proveedor en config sin modificar este módulo.
  - tenacity para retry/backoff: las APIs de LLM tienen rate limits y
    fallos transitorios. No queremos que un 429 mate el pipeline entero.
  - La extracción de texto es separada de la generación — más fácil de
    testear y reemplazar (ej: agregar OCR para PDFs escaneados).
  - El prompt está en español porque el contenido de la UNDEF es en español.

TODO: Una vez implementado el scraper real, ajustar MAX_TOKENS según el
      tamaño promedio de los documentos de Moodle de la UNDEF.
"""

import logging
import re
from dataclasses import dataclass
from pathlib import Path

from tenacity import (
    retry,
    retry_if_exception,
    retry_if_exception_type,
    stop_after_attempt,
    wait_exponential,
    before_sleep_log,
)

try:
    from scripts.academic_calendar import parse_llm_calendar_payload
except ImportError:  # pragma: no cover - fallback for direct script execution
    from academic_calendar import parse_llm_calendar_payload

log = logging.getLogger(__name__)

# Número máximo de intentos ante errores de rate limit o errores transitorios
MAX_RETRY_ATTEMPTS = 4

# Backoff exponencial: 2s, 4s, 8s, 16s entre reintentos
# Suficiente para sobrevivir rate limits típicos de APIs de LLM
RETRY_WAIT_MIN_SECONDS = 2
RETRY_WAIT_MAX_SECONDS = 30

# Máximo de tokens en el resumen generado.
# GitHub Models free tier (modelos low, ej: gpt-4o-mini): límite real = 4000 tokens/request.
# Se usa 3800 como margen de seguridad.
MAX_OUTPUT_TOKENS = 3800

# Límite de caracteres de entrada por proveedor.
# GitHub Models (gpt-4o-mini) tiene 8k tokens de input en el tier gratuito.
# Sistema + template del prompt consumen ~600 tokens → quedan ~7400 para contenido.
# 7400 tokens * 3.5 chars/token ≈ 25.000 chars seguros.
# Para documentos más largos habría que implementar chunking (dividir + merge).
# Gemini 1.5 Flash tiene contexto de 1M tokens — el límite es mucho más holgado.
MAX_INPUT_CHARS_GITHUB = 25_000
MAX_INPUT_CHARS_GEMINI = 80_000

# Límite de páginas para PDFs: documentos más largos se saltan
# para no exceder el contexto del LLM y evitar resúmenes de baja calidad
MAX_PDF_PAGES = 30


class DocumentTooLargeError(Exception):
    """El documento supera el límite de tamaño procesable por el pipeline."""


class DocumentUnreadableError(Exception):
    """El documento no tiene texto extraíble (PDF escaneado o basado en imágenes)."""


class DocumentTruncatedError(Exception):
    """El LLM truncó la respuesta antes de cerrar el bloque kdef-events.

    Ocurre cuando el documento genera un resumen más largo que MAX_OUTPUT_TOKENS.
    Incluye el markdown parcial generado antes del corte para que el pipeline
    pueda escribir un placeholder con el contenido disponible.
    """

    def __init__(self, partial_markdown: str = "") -> None:
        super().__init__("Respuesta truncada por el LLM antes de completar el bloque kdef-events")
        self.partial_markdown = partial_markdown


@dataclass
class DocumentSummary:
    """Resultado estructurado de la etapa de resumen."""

    markdown: str
    important_dates: list[dict[str, str]]

# Prompt del sistema — define el rol y las instrucciones de formato
SYSTEM_PROMPT = """Eres un asistente académico especializado en ciberseguridad y defensa nacional.
Tu tarea es generar resúmenes estructurados de materiales académicos para estudiantes
de la Licenciatura en Ciberdefensa de la UNDEF (Universidad de la Defensa Nacional, Argentina).

Reglas del resumen:
1. Escribir siempre en español rioplatense (Argentina)
2. Estructura: título, conceptos clave, desarrollo principal, conclusiones
3. Formato: markdown limpio, adecuado para Quartz (jardín de conocimiento)
4. Longitud: conciso pero completo — no cortar conceptos importantes
5. NO inventar información — si algo no está claro en el texto, decirlo explícitamente
6. NO expandir siglas ni acrónimos a menos que el propio documento los defina. Si una sigla
   aparece sin definición, usarla tal cual sin intentar interpretarla
7. Si el documento es una tabla, cronograma o lista de fechas: transcribir fielmente todos
   los eventos y fechas que contenga, sin parafrasear ni omitir filas
8. Incluir definiciones de términos técnicos solo cuando el documento las provea
9. Si el documento tiene ejercicios o preguntas, incluirlos resumidos
10. Al final de la respuesta devolver un bloque fenced con info string `kdef-events`
    y JSON válido con un campo `important_dates`
11. `important_dates` debe incluir TODOS los hitos académicos explícitos en el texto:
    finales, parciales, exámenes, entregas, recuperatorios, clases especiales u otros eventos.
    En documentos de tipo cronograma o calendario, cada fila o entrada es un evento separado.
    NO incluir fechas de publicación de artículos, papers o trabajos citados como referencia bibliográfica
12. Cada evento debe usar exclusivamente este schema:
    {"title": "...", "kind": "parcial|final|examen|entrega|clase|otro",
     "date_iso": "YYYY-MM-DD o null", "date_text": "texto original de la fecha",
     "time_text": "hora si existe, sino cadena vacía",
     "details": "contexto breve", "source_excerpt": "fragmento breve del material"}
13. Si no hay fechas, devolver `{"important_dates": []}` en ese bloque
14. No escribir nada después del bloque `kdef-events`"""

# Prompt del usuario — instrucción concreta sobre qué hacer con el texto
USER_PROMPT_TEMPLATE = """Generá un resumen académico estructurado en markdown del siguiente documento.

Documento: {filename}
Año académico en curso: {academic_year}

Contenido:
{content}

El resumen debe incluir:
- Una descripción de 2-3 oraciones de qué trata el documento
- Los conceptos y términos clave definidos brevemente
- Los puntos principales desarrollados en el documento
- Referencias o lecturas adicionales mencionadas (si las hay)

Formato final obligatorio:
1. Primero, el resumen visible en markdown
2. Después, un único bloque:
```kdef-events
{{"important_dates": [...]}}
```

Reglas para las fechas:
- Usar `date_iso` solo cuando el texto permita identificar día, mes y año con certeza
- Si la fecha es parcial o ambigua, dejar `date_iso` en null y completar `date_text`
- No inventar años ni horarios
- No duplicar eventos equivalentes dentro del mismo documento"""


def extract_text_from_pdf(path: Path) -> str:
    """
    Extraer texto de un archivo PDF.

    Args:
        path: Path al archivo PDF.

    Returns:
        Texto extraído como string.

    Raises:
        ValueError: Si el PDF está vacío o no se puede extraer texto.

    TODO: Considerar agregar fallback con OCR (pytesseract) para PDFs
          escaneados que no tienen texto seleccionable.
    """
    import pypdf

    text_parts: list[str] = []

    with open(path, "rb") as f:
        reader = pypdf.PdfReader(f)
        page_count = len(reader.pages)
        if page_count > MAX_PDF_PAGES:
            raise DocumentTooLargeError(
                f"{path.name} tiene {page_count} páginas (límite: {MAX_PDF_PAGES})"
            )
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"[Página {page_num + 1}]\n{page_text}")

    if not text_parts:
        raise DocumentUnreadableError(
            f"{path.name}: PDF sin texto extraíble — posiblemente escaneado o basado en imágenes"
        )

    full_text = "\n\n".join(text_parts)
    # Verificar que haya contenido real más allá de espacios y saltos de línea
    if len(full_text.replace(" ", "").replace("\n", "").replace("\t", "")) < 50:
        raise DocumentUnreadableError(
            f"{path.name}: PDF con texto insuficiente — posiblemente escaneado o basado en imágenes"
        )

    return full_text


def extract_text_from_docx(path: Path) -> str:
    """
    Extraer texto de un archivo DOCX.

    Args:
        path: Path al archivo DOCX.

    Returns:
        Texto extraído como string.

    Raises:
        ValueError: Si el documento está vacío o no tiene texto extraíble.
    """
    from docx import Document

    doc = Document(path)
    parts = [para.text for para in doc.paragraphs if para.text.strip()]

    if not parts:
        raise ValueError(f"No se pudo extraer texto del DOCX: {path.name}")

    return "\n".join(parts)


def extract_text_from_pptx(path: Path) -> str:
    """
    Extraer texto de un archivo PPTX.

    Recorre todas las diapositivas y extrae el texto de cada shape,
    etiquetando cada diapositiva para preservar la estructura.

    Args:
        path: Path al archivo PPTX.

    Returns:
        Texto extraído como string, con marcadores por diapositiva.

    Raises:
        ValueError: Si la presentación no tiene texto extraíble.
    """
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if slide_texts:
            parts.append(f"[Diapositiva {i}]\n" + "\n".join(slide_texts))

    if not parts:
        raise ValueError(f"No se pudo extraer texto del PPTX: {path.name}")

    return "\n\n".join(parts)


def extract_text(path: Path) -> str:
    """
    Extraer texto de un archivo según su extensión.

    Dispatcher que selecciona el extractor correcto según el tipo de archivo.

    Args:
        path: Path al archivo del que extraer texto.

    Returns:
        Texto extraído como string.

    Raises:
        ValueError: Si la extensión no está soportada.
    """
    ext = path.suffix.lower()

    extractors = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
        ".txt": lambda p: p.read_text(encoding="utf-8", errors="replace"),
        ".md": lambda p: p.read_text(encoding="utf-8", errors="replace"),
    }

    if ext not in extractors:
        raise ValueError(
            f"Extensión no soportada: {ext}. "
            f"Extensiones soportadas: {list(extractors.keys())}"
        )

    return extractors[ext](path)


def _summarize_with_github_models(text: str, filename: str, api_key: str, model_name: str, academic_year: str = "") -> str:
    """
    Generar resumen usando GitHub Models (API compatible con OpenAI).

    Args:
        text: Texto del documento a resumir.
        filename: Nombre del archivo (para contexto en el prompt).
        api_key: API key de GitHub Models.
        model_name: Nombre del modelo sin prefijo (ej: "gpt-4o-mini").
        academic_year: Año académico en curso (ej: "2026") para contextualizar fechas.

    Returns:
        Resumen en markdown generado por el LLM.
    """
    from openai import OpenAI

    client = OpenAI(
        api_key=api_key,
        base_url="https://models.inference.ai.azure.com",
    )

    response = client.chat.completions.create(
        model=model_name,
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {
                "role": "user",
                "content": USER_PROMPT_TEMPLATE.format(
                    filename=filename,
                    academic_year=academic_year or "desconocido",
                    content=text[:MAX_INPUT_CHARS_GITHUB],
                ),
            },
        ],
        max_tokens=MAX_OUTPUT_TOKENS,
        temperature=0.3,  # Baja temperatura para resúmenes más factuales
    )

    return response.choices[0].message.content or ""


def _summarize_with_openrouter(text: str, filename: str, api_key: str, model_name: str, academic_year: str = "") -> str:
    """
    Generar resumen usando OpenRouter (API compatible con OpenAI).

    OpenRouter expone modelos de distintos proveedores bajo una API unificada.
    Incluye modelos gratuitos con rate limit diario renovable.

    Args:
        text: Texto del documento a resumir.
        filename: Nombre del archivo (para contexto en el prompt).
        api_key: API key de OpenRouter.
        model_name: Nombre del modelo sin prefijo (ej: "meta-llama/llama-3.1-8b-instruct:free").
        academic_year: Año académico en curso (ej: "2026") para contextualizar fechas.

    Returns:
        Resumen en markdown generado por el LLM.
    """
    from openai import OpenAI

    client = OpenAI(
        api_key=api_key,
        base_url="https://openrouter.ai/api/v1",
        default_headers={
            "HTTP-Referer": "https://github.com/matzalazar/kdef",
            "X-Title": "kdef",
        },
    )

    # Fusionar system prompt en el mensaje de usuario: algunos modelos de OpenRouter
    # (ej: Gemma via Google AI Studio) no aceptan el rol "system".
    # El formato fusionado es universalmente compatible.
    user_content = (
        f"{SYSTEM_PROMPT}\n\n---\n\n"
        + USER_PROMPT_TEMPLATE.format(
            filename=filename,
            academic_year=academic_year or "desconocido",
            content=text[:MAX_INPUT_CHARS_GITHUB],
        )
    )

    response = client.chat.completions.create(
        model=model_name,
        messages=[{"role": "user", "content": user_content}],
        max_tokens=MAX_OUTPUT_TOKENS,
        temperature=0.3,
    )

    return response.choices[0].message.content or ""


def _summarize_with_gemini(text: str, filename: str, api_key: str, model_name: str, academic_year: str = "") -> str:
    """
    Generar resumen usando Google Gemini (fallback).

    Args:
        text: Texto del documento a resumir.
        filename: Nombre del archivo (para contexto en el prompt).
        api_key: API key de Google Gemini.
        model_name: Nombre del modelo sin prefijo (ej: "gemini-2.0-flash").
        academic_year: Año académico en curso (ej: "2026") para contextualizar fechas.

    Returns:
        Resumen en markdown generado por el LLM.
    """
    from google import genai
    from google.genai import types

    client = genai.Client(api_key=api_key)

    prompt = USER_PROMPT_TEMPLATE.format(
        filename=filename,
        academic_year=academic_year or "desconocido",
        content=text[:MAX_INPUT_CHARS_GEMINI],
    )

    response = client.models.generate_content(
        model=model_name,
        contents=prompt,
        config=types.GenerateContentConfig(
            system_instruction=SYSTEM_PROMPT,
            max_output_tokens=MAX_OUTPUT_TOKENS,
            temperature=0.3,
        ),
    )
    return response.text or ""


def _is_rate_limit_error(exc: BaseException) -> bool:
    """Detectar errores 429 de cualquier proveedor LLM."""
    # openai.RateLimitError cubre GitHub Models y OpenRouter
    try:
        from openai import RateLimitError as OpenAIRateLimitError
        if isinstance(exc, OpenAIRateLimitError):
            return True
    except ImportError:
        pass
    # google.genai puede lanzar excepciones distintas
    exc_str = type(exc).__name__.lower()
    return "ratelimit" in exc_str or "resourceexhausted" in exc_str


@retry(
    # Reintentar ante rate limit (429), errores transitorios del servidor (5xx)
    # y errores de red. No reintentar ante 400, 401, 403.
    retry=retry_if_exception_type((ConnectionError, TimeoutError)) | retry_if_exception(_is_rate_limit_error),
    stop=stop_after_attempt(MAX_RETRY_ATTEMPTS),
    wait=wait_exponential(
        multiplier=2,
        min=15,
        max=120,
    ),
    before_sleep=before_sleep_log(log, logging.WARNING),
    reraise=True,
)
def summarize_document(path: Path, model: str, academic_year: str = "") -> DocumentSummary:
    """
    Generar un resumen markdown estructurado de un documento usando un LLM.

    Esta es la función pública de este módulo. El decorador @retry maneja
    automáticamente los reintentos con backoff exponencial.

    La selección del proveedor LLM se hace por prefijo del nombre del modelo:
    - "github/..." → GitHub Models (API compatible con OpenAI)
    - "gemini/..." → Google Gemini

    Args:
        path: Path al archivo a resumir. Puede ser PDF, TXT, MD.
        model: Nombre del modelo con prefijo del proveedor.
                Ejemplos: "github/gpt-4o-mini", "gemini/gemini-1.5-flash"
        academic_year: Año académico en curso (ej: "2026"). Se incluye en el prompt
                       para que el LLM pueda inferir el año en fechas parciales como "21-may".

    Returns:
        DocumentSummary con el markdown visible y las fechas importantes detectadas.

    Raises:
        DocumentTruncatedError: Si el LLM cortó la respuesta antes de cerrar el bloque
                                kdef-events. Contiene el markdown parcial disponible.
        ValueError: Si el formato del archivo no está soportado o el
                    modelo especificado no tiene un proveedor reconocido.
        RuntimeError: Si todos los reintentos fallan.
    """
    import os

    log.info("Generando resumen de %s con modelo %s", path.name, model)

    # Extraer texto del documento
    text = extract_text(path)
    log.debug("Texto extraído: %d caracteres", len(text))

    # Seleccionar proveedor y generar resumen
    provider, _, model_name = model.partition("/")

    if provider == "github":
        api_key = os.getenv("GITHUB_MODELS_KEY", "")
        if not api_key:
            raise ValueError("GITHUB_MODELS_KEY no está configurado")
        raw_summary = _summarize_with_github_models(text, path.name, api_key, model_name, academic_year)

    elif provider == "openrouter":
        api_key = os.getenv("OPENROUTER_API_KEY", "")
        if not api_key:
            raise ValueError("OPENROUTER_API_KEY no está configurado")
        raw_summary = _summarize_with_openrouter(text, path.name, api_key, model_name, academic_year)

    elif provider == "gemini":
        api_key = os.getenv("GEMINI_API_KEY", "")
        if not api_key:
            raise ValueError("GEMINI_API_KEY no está configurado")
        raw_summary = _summarize_with_gemini(text, path.name, api_key, model_name, academic_year)

    elif model == "dry-run/mock":
        # Modelo mock para dry-run — retorna un resumen de ejemplo
        return DocumentSummary(
            markdown=f"""## Resumen de {path.name}

*[Dry-run mode — este resumen es un placeholder generado por el pipeline en modo simulación]*

Este archivo sería procesado por el LLM en una ejecución real del pipeline.
""",
            important_dates=[],
        )
    else:
        raise ValueError(
            f"Proveedor LLM no reconocido: '{provider}'. "
            f"Usar prefijo 'github/' o 'gemini/'."
        )

    # Detectar truncamiento: el bloque kdef-events fue abierto pero nunca cerrado.
    # Esto ocurre cuando el LLM alcanza MAX_OUTPUT_TOKENS a mitad del JSON.
    # No hay retry — es una limitación del tamaño del documento, no un error transitorio.
    fence_open = re.search(r"```\s*kdef-events", raw_summary, re.IGNORECASE)
    if fence_open:
        # Buscar el cierre del fence después de la apertura
        close_pos = raw_summary.find("```", fence_open.end())
        if close_pos == -1:
            partial = raw_summary[: fence_open.start()].strip()
            raise DocumentTruncatedError(partial)

    markdown, important_dates = parse_llm_calendar_payload(raw_summary)
    return DocumentSummary(markdown=markdown, important_dates=important_dates)
